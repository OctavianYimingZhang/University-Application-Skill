#!/usr/bin/env python3
from __future__ import annotations

import base64
import io
import json
import mimetypes
import re
import sys
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

MAX_FILE_BYTES = 25 * 1024 * 1024
MAX_BLOCK_CHARS = 5000
MAX_PREVIEW_CHARS = 12000
TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".html", ".htm", ".json"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


class HtmlTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(data)

    def text(self) -> str:
        return "\n".join(self.parts)


def clean_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = re.sub(r"[ \t\r\f\v]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def make_block(label: str, text: str) -> dict[str, str] | None:
    cleaned = clean_text(text)
    if not cleaned:
        return None
    if len(cleaned) > MAX_BLOCK_CHARS:
        cleaned = cleaned[:MAX_BLOCK_CHARS].rstrip() + "..."
    return {"label": label, "text": cleaned}


def extension_for(name: str, mime_type: str) -> str:
    suffix = Path(name).suffix.lower()
    if suffix:
        return suffix
    guessed = mimetypes.guess_extension(mime_type or "")
    return guessed.lower() if guessed else ""


def decode_payload(payload: dict[str, Any]) -> tuple[str, str, bytes]:
    name = str(payload.get("name") or payload.get("fileName") or "uploaded-file").strip() or "uploaded-file"
    mime_type = str(payload.get("mimeType") or mimetypes.guess_type(name)[0] or "application/octet-stream")
    content_base64 = payload.get("contentBase64")
    if not isinstance(content_base64, str) or not content_base64:
        raise ValueError("contentBase64 is required.")

    if "," in content_base64 and content_base64.lstrip().startswith("data:"):
        content_base64 = content_base64.split(",", 1)[1]
    try:
        data = base64.b64decode(content_base64, validate=True)
    except Exception as exc:  # noqa: BLE001
        raise ValueError("contentBase64 is not valid base64.") from exc
    if len(data) > MAX_FILE_BYTES:
        raise ValueError("File exceeds the 25MB extraction limit.")
    if not data:
        raise ValueError("File is empty.")
    return name, mime_type, data


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text_like(name: str, ext: str, data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    text = decode_text(data)
    if ext == ".json":
        try:
            text = json.dumps(json.loads(text), indent=2, ensure_ascii=False)
        except Exception:  # noqa: BLE001
            warnings.append("JSON could not be parsed; extracted raw text instead.")
    if ext in {".html", ".htm"}:
        parser = HtmlTextExtractor()
        try:
            parser.feed(text)
            text = parser.text()
        except Exception:  # noqa: BLE001
            warnings.append("HTML parsing failed; extracted raw text instead.")
    block = make_block("Document text", text)
    return [block] if block else []


def extract_pdf(data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    blocks: list[dict[str, str]] = []
    try:
        import pdfplumber  # type: ignore

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                block = make_block(f"Page {index}", page.extract_text() or "")
                if block:
                    blocks.append(block)
    except ImportError:
        warnings.append("pdfplumber is not installed; using pypdf fallback.")
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"pdfplumber extraction failed: {exc}")

    if blocks:
        return blocks

    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(io.BytesIO(data))
        if getattr(reader, "is_encrypted", False):
            warnings.append("PDF is encrypted; text extraction was not possible.")
            return blocks
        for index, page in enumerate(reader.pages, start=1):
            block = make_block(f"Page {index}", page.extract_text() or "")
            if block:
                blocks.append(block)
    except ImportError:
        warnings.append("pypdf is not installed; PDF fallback extraction unavailable.")
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"pypdf extraction failed: {exc}")

    if not blocks:
        warnings.append("No embedded PDF text was found. This may be a scanned PDF.")
    return blocks


def extract_docx(data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        warnings.append("python-docx is not installed; DOCX extraction unavailable.")
        return []

    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"DOCX could not be opened: {exc}")
        return []

    blocks: list[dict[str, str]] = []
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    block = make_block("Document paragraphs", "\n".join(paragraphs))
    if block:
        blocks.append(block)

    for table_index, table in enumerate(document.tables, start=1):
        lines: list[str] = []
        for row in table.rows:
            cells = [clean_text(cell.text) for cell in row.cells if clean_text(cell.text)]
            if cells:
                lines.append(" | ".join(cells))
        block = make_block(f"Table {table_index}", "\n".join(lines))
        if block:
            blocks.append(block)
    return blocks


def extract_pptx(data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        warnings.append("python-pptx is not installed; PPTX extraction unavailable.")
        return []

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"PPTX could not be opened: {exc}")
        return []

    blocks: list[dict[str, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [clean_text(cell.text) for cell in row.cells if clean_text(cell.text)]
                    if cells:
                        parts.append(" | ".join(cells))
        block = make_block(f"Slide {index}", "\n".join(parts))
        if block:
            blocks.append(block)
    return blocks


def extract_xlsx(data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        warnings.append("openpyxl is not installed; XLSX extraction unavailable.")
        return []

    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"XLSX could not be opened: {exc}")
        return []

    blocks: list[dict[str, str]] = []
    for sheet in workbook.worksheets:
        lines: list[str] = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = [str(value) for value in row if value not in (None, "")]
            if values:
                lines.append(f"R{row_index}: " + " | ".join(values))
            if len(lines) >= 80:
                warnings.append(f"{sheet.title} was truncated after 80 non-empty rows.")
                break
        block = make_block(f"Sheet {sheet.title}", "\n".join(lines))
        if block:
            blocks.append(block)
    return blocks


def register_image(data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    try:
        from PIL import Image  # type: ignore

        image = Image.open(io.BytesIO(data))
        warnings.append(f"Image registered ({image.width}x{image.height}). OCR is not configured in v1; add manual notes for writing inspiration.")
    except ImportError:
        warnings.append("Image registered. Pillow is not installed, and OCR is not configured in v1.")
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"Image could not be inspected: {exc}")
    return []


def extract_blocks(name: str, mime_type: str, data: bytes, warnings: list[str]) -> list[dict[str, str]]:
    ext = extension_for(name, mime_type)
    if ext in TEXT_EXTENSIONS or mime_type.startswith("text/"):
        return extract_text_like(name, ext, data, warnings)
    if ext == ".pdf" or mime_type == "application/pdf":
        return extract_pdf(data, warnings)
    if ext == ".docx":
        return extract_docx(data, warnings)
    if ext == ".pptx":
        return extract_pptx(data, warnings)
    if ext == ".xlsx":
        return extract_xlsx(data, warnings)
    if ext in IMAGE_EXTENSIONS or mime_type.startswith("image/"):
        return register_image(data, warnings)

    warnings.append(f"Unsupported file type for automatic extraction: {ext or mime_type}. Add manual notes in the Writing Studio.")
    return []


def preview_from_blocks(blocks: list[dict[str, str]]) -> str:
    text = "\n\n".join(f"{block['label']}: {block['text']}" for block in blocks)
    if len(text) > MAX_PREVIEW_CHARS:
        return text[:MAX_PREVIEW_CHARS].rstrip() + "..."
    return text


def main() -> None:
    try:
        payload = json.loads(sys.stdin.read())
        if not isinstance(payload, dict):
            raise ValueError("Request body must be a JSON object.")
        name, mime_type, data = decode_payload(payload)
        warnings: list[str] = []
        blocks = extract_blocks(name, mime_type, data, warnings)
        result = {
            "ok": bool(blocks) or any("Image registered" in warning for warning in warnings),
            "fileName": name,
            "mimeType": mime_type,
            "sizeBytes": len(data),
            "blocks": blocks,
            "preview": preview_from_blocks(blocks),
            "warnings": warnings,
        }
    except Exception as exc:  # noqa: BLE001
        result = {
            "ok": False,
            "fileName": "uploaded-file",
            "mimeType": "application/octet-stream",
            "sizeBytes": 0,
            "blocks": [],
            "preview": "",
            "warnings": [str(exc)],
            "error": str(exc),
        }

    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
